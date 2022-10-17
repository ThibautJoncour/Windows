#!/usr/bin/env python
# coding: utf-8

# In[9]:


from pptx import Presentation
from pptx.util import Inches
import pandas as pd
import os


# In[10]:


df=pd.read_excel(r'C:\Users\Mon PC\Documents\ppt_excel.xlsx')
df


# In[11]:


def df_to_table(slide,df,x, y, cx, cy,index_names=False,columns_names=True):
    df=pd.DataFrame(df)
    rows, cols = df.shape
    
    res = slide.shapes.add_table(rows+columns_names, cols+index_names, x, y, cx, cy) # NOMBRE COLONNE LIGNE + taille
    if columns_names:
        for col_index, col_name in enumerate(list(df.columns)):
            cell=res.table.cell(0,col_index+index_names)
            #cell.text_frame.fit_text(max_size=12)
            #cell.text_frame.text='%s'%(col_name)
            cell.text = '%s'%(col_name)
    if index_names:
        for col_index, col_name in enumerate(list(df.index)):
            cell=res.table.cell(col_index+columns_names,0)
            cell.text = '%s'%(col_name)
            #cell.text_frame.fit_text(max_size=12)
    m = df.values
    for row in range(rows):
        for col in range(cols):
            cell=res.table.cell(row+columns_names, col+index_names)
            if isinstance(m[row, col],float):
                cell.text = '%.2f'%(m[row, col])
            else:
                cell.text = '%s'%(m[row, col])
            #cell.text_frame.fit_text(max_size=12) 


# In[15]:


df=pd.read_excel(r'C:\Users\Mon PC\Documents\ppt_excel.xlsx')
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
x, y, cx, cy = Inches(0.01), Inches(0.2), Inches(10), Inches(0.1)

df_to_table(slide,df,x, y, cx, cy,index_names=False,columns_names=True)
prs.save(r'C:\Users\Mon PC\Documents\table2.pptx')
os.startfile(r'C:\Users\Mon PC\Documents\table2.pptx')


# In[ ]:




